import json
import logging
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def create_carousel_presentation(
        topic: str,
        slides_content: List[Dict],
        template_path: str = "templates/default.pptx") -> Presentation:
    """Create a PowerPoint presentation using a template and generated content."""
    if not slides_content:
        logger.error("No slides content provided")
        raise ValueError("No slides content provided")

    logger.info(f"Creating presentation with {len(slides_content)} slides")
    logger.info(
        f"First slide content: {json.dumps(slides_content[0], indent=2)}")
    """Create a PowerPoint presentation using a template and generated content."""
    try:
        logger.info(f"Opening template from: {template_path}")
        prs = Presentation(template_path)

        # Create placeholder dictionary
        placeholders = {}
        for i, slide_content in enumerate(slides_content):
            # Add title placeholder
            title_key = f"{{{{Title{i}}}}}"
            title_value = slide_content.get('title', '')
            placeholders[title_key] = title_value
            logger.info(
                f"Adding title placeholder: {title_key} -> {title_value}")

            # Add body placeholder with bullet points
            body_key = f"{{{{Body{i}}}}}"
            points = slide_content.get('points', [])
            body_value = "• " + "\n• ".join(points)
            placeholders[body_key] = body_value
            logger.info(
                f"Adding body placeholder: {body_key} -> {len(points)} bullet points"
            )

        # Process each slide
        for slide_idx, slide in enumerate(prs.slides):
            logger.info(f"Processing slide {slide_idx + 1}")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # Process each paragraph
                for paragraph in shape.text_frame.paragraphs:
                    # Process each run in the paragraph
                    for run in paragraph.runs:
                        original_text = run.text
                        modified_text = original_text

                        # Replace all placeholders in this run
                        for placeholder, replacement in placeholders.items():
                            if placeholder in modified_text:
                                logger.info(
                                    f"Found placeholder {placeholder} in text: {modified_text}"
                                )
                                modified_text = modified_text.replace(
                                    placeholder, replacement)
                                logger.info(f"Replaced with: {replacement}")

                        # Only update if text was modified
                        if modified_text != original_text:
                            run.text = modified_text

        logger.info("Successfully created presentation from template")
        return prs
    except Exception as e:
        logger.error(f"Failed to create presentation from template: {str(e)}")
        raise Exception(
            f"Failed to create presentation from template: {str(e)}")
