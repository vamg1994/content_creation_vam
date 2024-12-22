import os
import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import List

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

TEMPLATES_DIR = "templates"
DEFAULT_TEMPLATES = ["LinkedIn_Carrousel_VAM", "Linkedin_carrousel_vam2"]

def ensure_template_directory():
    """Ensure templates directory exists."""
    if not os.path.exists(TEMPLATES_DIR):
        os.makedirs(TEMPLATES_DIR)
        logger.info(f"Created templates directory: {TEMPLATES_DIR}")

def create_default_template(template_name: str) -> None:
    """Create a default template with the specified name."""
    try:
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Set slide dimensions and properties based on template style
        if template_name == "modern":
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
        elif template_name == "professional":
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
        elif template_name == "creative":
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
        # Add content placeholder slides
        for i in range(2):  # Add 2 content slides
            content_layout = prs.slide_layouts[1]  # Using a layout with title and content
            content_slide = prs.slides.add_slide(content_layout)
            content_title = content_slide.shapes.title
            content_body = content_slide.placeholders[1]
            
            # Add placeholder text
            content_title.text = f"{{Title{i}}}"
            content_body.text = f"{{Body{i}}}"
        
        # Save template
        template_path = os.path.join(TEMPLATES_DIR, f"{template_name}.pptx")
        prs.save(template_path)
        logger.info(f"Created template: {template_path}")
    except Exception as e:
        logger.error(f"Failed to create template {template_name}: {str(e)}")
        raise

def initialize_templates():
    """Initialize default templates if they don't exist."""
    try:
        ensure_template_directory()
        
        templates_created = False
        templates_status = []
        
        for template_name in DEFAULT_TEMPLATES:
            template_path = os.path.join(TEMPLATES_DIR, f"{template_name}.pptx")
            if not os.path.exists(template_path):
                try:
                    create_default_template(template_name)
                    templates_created = True
                    templates_status.append(f"Created template: {template_name}")
                except Exception as e:
                    error_msg = f"Failed to create template {template_name}: {str(e)}"
                    logger.error(error_msg)
                    templates_status.append(error_msg)
            else:
                templates_status.append(f"Template exists: {template_name}")
        
        if templates_created:
            logger.info("Successfully initialized new templates")
        
        logger.info("Templates status:\n" + "\n".join(templates_status))
        return templates_status
            
    except Exception as e:
        error_msg = f"Failed to initialize templates: {str(e)}"
        logger.error(error_msg)
        raise Exception(error_msg)

def get_available_templates() -> List[str]:
    """Get list of available templates."""
    try:
        ensure_template_directory()
        
        # Initialize templates if directory is empty
        if not os.path.exists(TEMPLATES_DIR) or not os.listdir(TEMPLATES_DIR):
            logger.info("No templates found. Initializing default templates...")
            initialize_templates()
        
        templates = []
        for template_name in DEFAULT_TEMPLATES:
            template_path = os.path.join(TEMPLATES_DIR, f"{template_name}.pptx")
            if os.path.exists(template_path):
                templates.append(template_name.capitalize())
        
        if not templates:
            logger.warning("No templates available. Creating default template...")
            create_default_template("default")
            templates = ["Default"]
            
        return sorted(templates)
    except Exception as e:
        logger.error(f"Failed to get available templates: {str(e)}")
        return ["Default"]
